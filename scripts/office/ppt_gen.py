#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT 生成器 — 接收结构化 JSON 描述，用 python-pptx 渲染 .pptx 文件。

用法:
  python ppt_gen.py <input.json> <output_dir>

input.json 结构:
{
  "title": "PPT 标题",
  "subtitle": "副标题（可选）",
  "author": "作者（可选）",
  "theme": "blue|green|purple|dark|orange（可选，默认 blue）",
  "sections": [
    { "title": "章节标题", "bullets": ["要点1", "要点2"], "notes": "演讲备注（可选）" }
  ]
}

输出: 打印 JSON 到 stdout:
  {"success": true, "path": "D:/DesktopPetData/Documents/xxx.pptx", "title": "..."}
  {"success": false, "error": "..."}
"""

import json
import sys
import os
import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── 主题色 ─────────────────────────────────────────────────────────────────
THEMES = {
    "blue":   {"primary": RGBColor(0x1F, 0x4E, 0x79), "accent": RGBColor(0x2E, 0x86, 0xC1), "light": RGBColor(0xDE, 0xEB, 0xF7)},
    "green":  {"primary": RGBColor(0x1E, 0x6B, 0x3C), "accent": RGBColor(0x33, 0x8A, 0x5A), "light": RGBColor(0xDD, 0xF0, 0xE2)},
    "purple": {"primary": RGBColor(0x5B, 0x2C, 0x6F), "accent": RGBColor(0x7B, 0x4F, 0x9A), "light": RGBColor(0xE9, 0xDE, 0xF1)},
    "dark":   {"primary": RGBColor(0x2D, 0x2D, 0x2D), "accent": RGBColor(0x8A, 0x8A, 0x8A), "light": RGBColor(0xE8, 0xE8, 0xE8)},
    "orange": {"primary": RGBColor(0xB4, 0x50, 0x0E), "accent": RGBColor(0xE6, 0x7E, 0x22), "light": RGBColor(0xFB, 0xE9, 0xD7)},
}
FONT_CN = "微软雅黑"

def sanitize_filename(name):
    """清理非法文件名字符"""
    return "".join(c for c in name if c not in '<>:"/\\|?*').strip() or "presentation"


def add_bg_rect(slide, color, left, top, width, height):
    """添加纯色矩形背景"""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False
    return rect


def add_textbox(slide, left, top, width, height, text, size, color, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = FONT_CN
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, bullets, size, color, line_spacing=1.25):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(8)
        try:
            p.line_spacing = line_spacing
        except Exception:
            pass
        # 项目符号前缀
        marker = p.add_run()
        marker.text = "▸ "
        marker.font.name = FONT_CN
        marker.font.size = Pt(size)
        marker.font.bold = True
        marker.font.color.rgb = RGBColor(0x2E, 0x86, 0xC1)
        run = p.add_run()
        run.text = b
        run.font.name = FONT_CN
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def build_cover(prs, theme, title, subtitle, author):
    """封面页"""
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    sw, sh = prs.slide_width, prs.slide_height
    # 背景
    add_bg_rect(s, theme["primary"], 0, 0, sw, sh)
    # 装饰条
    add_bg_rect(s, theme["accent"], 0, sh - Inches(0.18), sw, Inches(0.18))
    # 标题
    add_textbox(s, Inches(0.9), Inches(2.2), sw - Inches(1.8), Inches(1.4),
                title, 40, RGBColor(0xFF, 0xFF, 0xFF), bold=True, align=PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(s, Inches(1.5), Inches(3.7), sw - Inches(3.0), Inches(0.8),
                    subtitle, 18, theme["light"], align=PP_ALIGN.CENTER)
    if author:
        add_textbox(s, Inches(1.5), Inches(4.9), sw - Inches(3.0), Inches(0.6),
                    author, 14, theme["light"], align=PP_ALIGN.CENTER)
    date_str = datetime.date.today().strftime("%Y-%m-%d")
    add_textbox(s, Inches(1.5), Inches(5.5), sw - Inches(3.0), Inches(0.5),
                date_str, 12, theme["light"], align=PP_ALIGN.CENTER)


def build_section(prs, theme, idx, total, section):
    """内容页"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height
    # 背景 + 顶部色带
    add_bg_rect(s, RGBColor(0xFF, 0xFF, 0xFF), 0, 0, sw, sh)
    add_bg_rect(s, theme["primary"], 0, 0, sw, Inches(0.12))
    # 页码
    add_textbox(s, sw - Inches(1.2), sh - Inches(0.5), Inches(1.0), Inches(0.4),
                f"{idx} / {total}", 12, RGBColor(0x99, 0x99, 0x99), align=PP_ALIGN.RIGHT)
    # 章节标题
    add_textbox(s, Inches(0.7), Inches(0.5), sw - Inches(1.4), Inches(0.9),
                section.get("title", f"第 {idx} 节"), 28, theme["primary"], bold=True)
    # 分隔线
    add_bg_rect(s, theme["accent"], Inches(0.7), Inches(1.35), Inches(1.2), Inches(0.06))
    # 要点
    bullets = section.get("bullets") or []
    if bullets:
        add_bullets(s, Inches(0.9), Inches(1.8), sw - Inches(1.8), sh - Inches(2.6),
                    bullets, 18, RGBColor(0x33, 0x33, 0x33))
    # 备注
    notes = section.get("notes")
    if notes:
        s.notes_slide.notes_text_frame.text = notes


def build_closing(prs, theme, title):
    """结束页"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height
    add_bg_rect(s, theme["primary"], 0, 0, sw, sh)
    add_textbox(s, Inches(0.9), Inches(3.0), sw - Inches(1.8), Inches(1.2),
                "谢谢观看", 44, RGBColor(0xFF, 0xFF, 0xFF), bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.9), Inches(4.3), sw - Inches(1.8), Inches(0.6),
                title, 16, theme["light"], align=PP_ALIGN.CENTER)


def main():
    if len(sys.argv) < 3:
        print(json.dumps({"success": False, "error": "用法: python ppt_gen.py <input.json> <output_dir>"}, ensure_ascii=False))
        return 1

    input_path = sys.argv[1]
    out_dir = sys.argv[2]

    try:
        # utf-8-sig 兼容 PowerShell 写入的 BOM，也兼容 Node 无 BOM 写入
        with open(input_path, "r", encoding="utf-8-sig") as f:
            desc = json.load(f)
    except Exception as e:
        print(json.dumps({"success": False, "error": f"读取描述文件失败: {e}"}, ensure_ascii=False))
        return 1

    try:
        title = str(desc.get("title") or "演示文稿")
        subtitle = str(desc.get("subtitle") or "")
        author = str(desc.get("author") or "")
        theme_name = str(desc.get("theme") or "blue").lower()
        theme = THEMES.get(theme_name, THEMES["blue"])
        sections = desc.get("sections") or []

        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9
        prs.slide_height = Inches(7.5)

        build_cover(prs, theme, title, subtitle, author)
        total = max(len(sections), 1)
        for i, sec in enumerate(sections, start=1):
            build_section(prs, theme, i, total, sec)
        build_closing(prs, theme, title)

        os.makedirs(out_dir, exist_ok=True)
        safe_title = sanitize_filename(title)
        filename = f"{safe_title}_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        path = os.path.join(out_dir, filename)
        prs.save(path)

        print(json.dumps({
            "success": True,
            "path": path.replace("\\", "/"),
            "title": title,
            "slides": len(sections) + 2
        }, ensure_ascii=False))
        return 0
    except Exception as e:
        print(json.dumps({"success": False, "error": f"生成失败: {e}"}, ensure_ascii=False))
        return 1


if __name__ == "__main__":
    sys.exit(main())
